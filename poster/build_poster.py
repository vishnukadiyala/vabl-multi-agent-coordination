"""Build VABL DISC poster from Slide 4 template."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

prs = Presentation('poster/Poster Templates.pptx')
slide = prs.slides[3]  # Slide 4 template

# Update title bar
for shape in slide.shapes:
    if not shape.has_text_frame:
        continue
    text = shape.text_frame.text

    if 'Project Title' in text:
        shape.text_frame.clear()
        p = shape.text_frame.paragraphs[0]
        run = p.add_run()
        run.text = "Implicit Coordination via Attention-Driven Latent Belief Representations"
        run.font.size = Pt(28)
        run.font.bold = True
        run.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER

        p2 = shape.text_frame.add_paragraph()
        run2 = p2.add_run()
        run2.text = "Vishnu Kadiyala, Mohammed Atiquzzaman \u2014 School of Computer Science, University of Oklahoma"
        run2.font.size = Pt(16)
        run2.font.color.rgb = RGBColor(220, 220, 220)
        p2.alignment = PP_ALIGN.CENTER

    elif 'Logo' in text:
        shape.text_frame.clear()
        p = shape.text_frame.paragraphs[0]
        run = p.add_run()
        run.text = "OU"
        run.font.size = Pt(24)
        run.font.bold = True
        run.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER

    elif 'Acknowledgement' in text:
        shape.text_frame.clear()
        p = shape.text_frame.paragraphs[0]
        run = p.add_run()
        run.text = "Contact: vishnupk@ou.edu  |  github.com/vishnukadiyala/vabl-multi-agent-coordination"
        run.font.size = Pt(12)
        run.font.color.rgb = RGBColor(100, 100, 100)


def add_text(slide, left, top, width, height, lines, default_size=14):
    """Add a text box with formatted lines. Lines are (text, size, bold, color) tuples or strings."""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, line in enumerate(lines):
        if isinstance(line, str):
            line = (line, default_size, False, RGBColor(40, 40, 40))
        text, size, bold, color = line

        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        if text == "":
            p.space_after = Pt(4)
            continue

        run = p.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color

B = True
TITLE_COLOR = RGBColor(0, 51, 102)
BODY_COLOR = RGBColor(40, 40, 40)
SUB_COLOR = RGBColor(80, 80, 80)
GREEN = RGBColor(0, 100, 0)
GRAY = RGBColor(100, 100, 100)

# ============================================================
# Panel 1: Motivation (top-left, ~11.4x9.6 at 0.7, 3.3)
# ============================================================
add_text(slide, 1.0, 3.5, 10.5, 9.0, [
    ("The Problem: Coordination Collapse", 24, B, TITLE_COLOR),
    ("", 10, False, BODY_COLOR),
    ("Multi-agent coordination under partial observability", 15, False, BODY_COLOR),
    ("requires reasoning about teammate intent. Standard", 15, False, BODY_COLOR),
    ("methods use egocentric memory that fails to model", 15, False, BODY_COLOR),
    ("teammate behavior, leading to coordination collapse.", 15, False, BODY_COLOR),
    ("", 10, False, BODY_COLOR),
    ("The collapse phenomenon:", 16, B, BODY_COLOR),
    ("\u2022 Agents discover coordination but cannot sustain it", 14, False, SUB_COLOR),
    ("\u2022 MAPPO: 100% collapse at 10M steps (peak 503, final 0)", 14, False, SUB_COLOR),
    ("\u2022 AERIAL (Phan et al., 2023): also 100% collapse", 14, False, SUB_COLOR),
    ("\u2022 TarMAC (communication): 85% collapse (final ~54)", 14, False, SUB_COLOR),
    ("", 10, False, BODY_COLOR),
    ("Our answer: VABL", 16, B, GREEN),
    ("Attention over teammate actions + auxiliary", 15, False, BODY_COLOR),
    ("prediction loss = stable implicit coordination", 15, False, BODY_COLOR),
    ("without any communication channel.", 15, B, BODY_COLOR),
])

# ============================================================
# Panel 2: Architecture + Results (top-right WIDE, ~22.9x9.6 at 12.3, 3.3)
# ============================================================
add_text(slide, 12.6, 3.5, 22.0, 0.6, [
    ("Architecture & Main Results", 24, B, TITLE_COLOR),
])

# Architecture diagram
slide.shapes.add_picture('poster/fig_architecture.png',
    Inches(12.6), Inches(4.2), Inches(10.8), Inches(4.3))

# Collapse comparison
slide.shapes.add_picture('poster/fig_collapse_comparison.png',
    Inches(23.6), Inches(4.2), Inches(11.0), Inches(4.3))

# Key numbers banner
add_text(slide, 12.6, 8.8, 22.0, 3.5, [
    ("3.3\u00d7 higher reward than MAPPO on Overcooked  \u2502  12\u00d7 lower variance  \u2502  38% collapse vs 100% (MAPPO/AERIAL)", 16, B, GREEN),
    ("", 8, False, BODY_COLOR),
    ("VABL trades ~48% of TarMAC\u2019s reward for eliminating communication \u2014 favorable when bandwidth is constrained.", 14, False, SUB_COLOR),
])

# ============================================================
# Panel 3: Method (bottom-left, ~11.4x9.6 at 0.7, 13.1)
# ============================================================
add_text(slide, 1.0, 13.3, 10.5, 9.0, [
    ("How VABL Works", 24, B, TITLE_COLOR),
    ("", 10, False, BODY_COLOR),
    ("Each agent maintains a belief state updated by:", 15, False, BODY_COLOR),
    ("", 8, False, BODY_COLOR),
    ("1. Attention over Teammate Actions", 16, B, BODY_COLOR),
    ("   Multi-head attention aggregates visible teammate", 14, False, SUB_COLOR),
    ("   actions. Query = belief, Keys/Values = actions.", 14, False, SUB_COLOR),
    ("", 8, False, BODY_COLOR),
    ("2. GRU Belief Update", 16, B, BODY_COLOR),
    ("   Observation encoding + attention context feed", 14, False, SUB_COLOR),
    ("   into recurrent belief state.", 14, False, SUB_COLOR),
    ("", 8, False, BODY_COLOR),
    ("3. Auxiliary Prediction (Theory of Mind)", 16, B, BODY_COLOR),
    ("   Predicts teammates\u2019 next actions from belief.", 14, False, SUB_COLOR),
    ("   Regularizes beliefs, preventing collapse.", 14, False, SUB_COLOR),
    ("   Accuracy: 17% (chance) \u2192 86% during training.", 14, False, SUB_COLOR),
    ("", 8, False, BODY_COLOR),
    ("No communication channel required.", 16, B, GREEN),
    ("< 10% overhead vs MAPPO (~1.3\u00d7 parameters).", 14, False, SUB_COLOR),
])

# ============================================================
# Panel 4: Ablation (bottom-middle, ~11.4x9.6 at 12.3, 13.1)
# ============================================================
add_text(slide, 12.6, 13.3, 10.5, 0.6, [
    ("Component Ablation", 22, B, TITLE_COLOR),
])

slide.shapes.add_picture('poster/fig_ablation_cramped.png',
    Inches(12.6), Inches(14.1), Inches(10.5), Inches(7.5))

add_text(slide, 12.6, 21.8, 10.5, 1.0, [
    ("Aux loss: consistent variance reduction (2.3\u00d7) across layouts.", 13, False, SUB_COLOR),
    ("Attention: helps at N>2, layout-dependent at N=2.", 13, False, SUB_COLOR),
])

# ============================================================
# Panel 5: Key Insight + Future (bottom-right, ~11.4x9.6 at 23.9, 13.1)
# ============================================================
add_text(slide, 24.2, 13.3, 10.5, 0.6, [
    ("Discovery vs. Maintenance", 22, B, TITLE_COLOR),
])

slide.shapes.add_picture('poster/fig_discovery_vs_maintenance.png',
    Inches(24.2), Inches(14.1), Inches(10.5), Inches(5.5))

add_text(slide, 24.2, 19.8, 10.5, 3.0, [
    ("Key Insight", 18, B, TITLE_COLOR),
    ("Every method can discover coordination.", 14, False, BODY_COLOR),
    ("The challenge is keeping it. VABL\u2019s auxiliary", 14, False, BODY_COLOR),
    ("prediction loss prevents belief drift, sustaining", 14, False, BODY_COLOR),
    ("coordination without communication.", 14, B, GREEN),
    ("", 8, False, BODY_COLOR),
    ("Future Work", 16, B, TITLE_COLOR),
    ("\u2022 Hanabi evaluation (genuine partial observability)", 13, False, SUB_COLOR),
    ("\u2022 Scaling to N>10 agents (sparse attention)", 13, False, SUB_COLOR),
    ("\u2022 Hybrid implicit + explicit communication", 13, False, SUB_COLOR),
])

# Delete other template slides (keep only slide 4 = index 3)
for idx in sorted([0, 1, 2, 4], reverse=True):
    rId = prs.slides._sldIdLst[idx].rId
    prs.part.drop_rel(rId)
    del prs.slides._sldIdLst[idx]

prs.save('poster/VABL_Poster.pptx')
print("Poster saved: poster/VABL_Poster.pptx")

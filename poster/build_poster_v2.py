"""Build VABL DISC poster v2 — competition-optimized.

Layout (Slide 4 template, 36x24 inches):
+------------------------------------------------------------------+
| TITLE: Learning to Coordinate Without Talking                     |
+-------------+----------------------------------------------------+
| 1. PROBLEM  |  2. HOW VABL WORKS (wide)                          |
|   + why it  |    Architecture + Hero figure (Reward Retained %)   |
|   matters   |    + TL;DR callout                                  |
+-------------+------------------------+---------------------------+
| 3. EVIDENCE                          | 4. CONCLUSIONS & IMPACT   |
|   Learning curve + Variance + Table  |   Takeaways + Applications|
+--------------------------------------+---------------------------+
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

B = True
TITLE_COLOR = RGBColor(0, 51, 102)
BODY_COLOR = RGBColor(40, 40, 40)
SUB_COLOR = RGBColor(80, 80, 80)
GREEN = RGBColor(0, 107, 63)
WHITE = RGBColor(255, 255, 255)
GRAY = RGBColor(100, 100, 100)
HIGHLIGHT_BG = RGBColor(0, 107, 63)

prs = Presentation('poster/Poster Templates.pptx')
slide = prs.slides[3]  # Slide 4 template


def add_text(slide, left, top, width, height, lines, default_size=14):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if isinstance(line, str):
            line = (line, default_size, False, BODY_COLOR)
        text, size, bold, color = line
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if text == "":
            p.space_after = Pt(4)
            continue
        run = p.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color


# ============================================================
# TITLE BAR
# ============================================================
for shape in slide.shapes:
    if not shape.has_text_frame:
        continue
    text = shape.text_frame.text
    if 'Project Title' in text:
        shape.text_frame.clear()
        p = shape.text_frame.paragraphs[0]
        run = p.add_run()
        run.text = "Learning to Coordinate Without Talking:"
        run.font.size = Pt(30)
        run.font.bold = True
        run.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

        p2 = shape.text_frame.add_paragraph()
        run2 = p2.add_run()
        run2.text = "Data-Driven Belief Models for Robust Multi-Agent Coordination"
        run2.font.size = Pt(22)
        run2.font.color.rgb = RGBColor(200, 230, 200)
        p2.alignment = PP_ALIGN.CENTER

        p3 = shape.text_frame.add_paragraph()
        run3 = p3.add_run()
        run3.text = "Vishnu Kadiyala, Mohammed Atiquzzaman \u2014 School of Computer Science, University of Oklahoma"
        run3.font.size = Pt(14)
        run3.font.color.rgb = RGBColor(200, 200, 200)
        p3.alignment = PP_ALIGN.CENTER

    elif 'Logo' in text:
        shape.text_frame.clear()
        p = shape.text_frame.paragraphs[0]
        run = p.add_run()
        run.text = "OU"
        run.font.size = Pt(22)
        run.font.bold = True
        run.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

    elif 'Acknowledgement' in text:
        shape.text_frame.clear()
        p = shape.text_frame.paragraphs[0]
        run = p.add_run()
        run.text = "Contact: vishnupk@ou.edu  \u2502  github.com/vishnukadiyala/vabl-multi-agent-coordination  \u2502  DISC 2026"
        run.font.size = Pt(11)
        run.font.color.rgb = GRAY


# ============================================================
# PANEL 1: THE PROBLEM (top-left, 11.4x9.6 at 0.7, 3.3)
# ============================================================
add_text(slide, 1.0, 3.5, 10.5, 9.0, [
    ("\u2460 The Problem", 22, B, TITLE_COLOR),
    ("", 4, False, BODY_COLOR),
    ("Can AI agents learn to work together", 16, B, BODY_COLOR),
    ("without talking to each other?", 16, B, BODY_COLOR),
    ("", 4, False, BODY_COLOR),
    ("A QB reads the safety cheat left, sees his receiver", 13, False, BODY_COLOR),
    ("break right, and throws \u2014 no audible needed. He", 13, False, BODY_COLOR),
    ("coordinates by predicting, not communicating.", 13, False, BODY_COLOR),
    ("", 4, False, BODY_COLOR),
    ("Setting: Decentralized POMDPs (Dec-POMDPs) \u2014", 13, B, BODY_COLOR),
    ("multiple agents act on partial observations with", 13, False, BODY_COLOR),
    ("no shared communication channel.", 13, False, BODY_COLOR),
    ("", 6, False, BODY_COLOR),
    ("The hidden problem: Coordination Collapse", 14, B, RGBColor(192, 57, 43)),
    ("Agents learn coordination, then catastrophically", 12, False, SUB_COLOR),
    ("forget it \u2014 like a team that falls apart in the 4th:", 12, False, SUB_COLOR),
    ("\u2022 MAPPO (PPO + centralized critic): 100% collapse", 12, False, SUB_COLOR),
    ("\u2022 AERIAL (attention, no aux loss): 100% collapse", 12, False, SUB_COLOR),
    ("\u2022 TarMAC (explicit communication): 85% collapse", 12, False, SUB_COLOR),
    ("", 6, False, BODY_COLOR),
    ("Our solution: VABL \u2014 agents build predictive", 13, B, GREEN),
    ("belief models of teammates from observable", 13, B, GREEN),
    ("actions via attention + auxiliary prediction.", 13, B, GREEN),
    ("Only 38% collapse, no communication required.", 13, B, GREEN),
])

# Collapse figure — fill remaining space in panel 1
slide.shapes.add_picture('poster/fig_collapse_all_methods.png',
    Inches(0.9), Inches(9.0), Inches(11.0), Inches(3.6))


# ============================================================
# PANEL 2: HOW VABL WORKS (top-right WIDE, 22.9x9.6 at 12.3, 3.3)
# ============================================================
add_text(slide, 12.6, 3.5, 22.0, 0.6, [
    ("\u2461 How VABL Works \u2014 Data-Driven Belief Learning", 22, B, TITLE_COLOR),
])

# Architecture diagram (left half of wide panel — fill more space)
slide.shapes.add_picture('poster/fig_architecture_v2.png',
    Inches(12.4), Inches(4.2), Inches(12.0), Inches(4.8))

# Hero figure (right half of wide panel — fill to panel edge)
slide.shapes.add_picture('poster/fig_hero_retained.png',
    Inches(24.2), Inches(4.0), Inches(11.0), Inches(6.2))

# Method steps (below architecture — shift up slightly)
add_text(slide, 12.4, 9.2, 11.5, 3.5, [
    ("CTDE paradigm: centralized critic (training), decentralized actors (execution)", 12, False, GRAY),
    ("", 4, False, BODY_COLOR),
    ("\u2776 Observe \u2014 Encode teammate actions via MLP", 13, False, SUB_COLOR),
    ("    (one-hot \u2192 d_e=64 embedding per teammate)", 11, False, GRAY),
    ("\u2777 Attend \u2014 Multi-head attention (h=4) over teammate", 13, False, SUB_COLOR),
    ("    action embeddings. Query = belief state b_t,", 11, False, GRAY),
    ("    Keys/Values = encoded actions. Produces social", 11, False, GRAY),
    ("    context vector c_t (permutation-invariant).", 11, False, GRAY),
    ("\u2778 Update \u2014 GRU cell: b_t = GRU([obs_enc || c_t], b_{t-1})", 13, False, SUB_COLOR),
    ("\u2779 Predict \u2014 Auxiliary MLP predicts teammates' next", 13, False, SUB_COLOR),
    ("    actions from b_t. Maximizes variational MI bound", 11, False, GRAY),
    ("    I(b_t; a_{t+1}). Keeps beliefs calibrated.", 11, False, GRAY),
    ("", 4, False, BODY_COLOR),
    ("L_total = L_PPO + 0.05 * L_aux_prediction", 12, B, BODY_COLOR),
    ("Aux accuracy: 17% (chance) \u2192 86%. Overhead: <10%", 12, False, GRAY),
])

# TL;DR callout box — shift up to fill gap
tldr = slide.shapes.add_textbox(Inches(24.2), Inches(10.5), Inches(11.0), Inches(2.2))
tf = tldr.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
run = p.add_run()
run.text = "VABL retains 62% of coordination vs 0% (MAPPO/AERIAL) \u2014 without any communication channel."
run.font.size = Pt(14)
run.font.bold = True
run.font.color.rgb = GREEN
p.alignment = PP_ALIGN.CENTER


# ============================================================
# PANEL 3: EVIDENCE (bottom-left, 11.4x9.6 at 0.7, 13.1 — extended to span 2)
# ============================================================
add_text(slide, 1.0, 13.3, 22.0, 0.6, [
    ("\u2462 Evidence: Rigorous Evaluation Across Environments", 22, B, TITLE_COLOR),
])

# Variance figure — expand to fill
slide.shapes.add_picture('poster/fig_variance_v2.png',
    Inches(0.9), Inches(14.0), Inches(7.8), Inches(5.8))

# Learning curves (Overcooked) — expand
slide.shapes.add_picture('figures/comparison_overcooked_asymmetric_advantages_updated.png',
    Inches(8.9), Inches(14.0), Inches(7.0), Inches(5.8))

# Ablation (Cramped Room) — expand
slide.shapes.add_picture('poster/fig_ablation_cramped.png',
    Inches(16.1), Inches(14.0), Inches(7.2), Inches(5.8))

# Results summary text below figures — tighter
add_text(slide, 0.9, 19.9, 22.5, 2.8, [
    ("Fig A: VABL 12\u00d7 more reliable", 12, B, GREEN),
    ("  across seeds (std 1.7 vs 20.3).", 11, False, SUB_COLOR),
    ("Fig B: VABL (green) sustains reward;", 12, B, GREEN),
    ("  MAPPO (blue) collapses after ep 200.", 11, False, SUB_COLOR),
    ("Fig C: Removing components degrades", 12, B, GREEN),
    ("  peak \u2014 full VABL is most reliable.", 11, False, SUB_COLOR),
    ("", 4, False, BODY_COLOR),
    ("Summary (5 seeds, 95% CIs):", 13, B, BODY_COLOR),
    ("\u2022 Overcooked AA: VABL 3.3\u00d7 MAPPO final reward, 38% vs 76% collapse", 11, False, SUB_COLOR),
    ("\u2022 5-agent Simple Coord: Best 95.7\u00b13.3 vs MAPPO 84.0\u00b110.4", 11, False, SUB_COLOR),
    ("\u2022 Cramped Room: Full VABL Best 1030\u00b170 (lowest variance)", 11, False, SUB_COLOR),
    ("\u2022 10M steps: MAPPO & AERIAL collapse to 0; TarMAC to 54", 11, False, SUB_COLOR),
    ("Environments: Overcooked-AI (2 agents), Simple Coordination (3-5 agents, stochastic visibility)", 10, False, GRAY),
])


# ============================================================
# PANEL 4: CONCLUSIONS & IMPACT (bottom-right, 11.4x9.6 at 23.9, 13.1)
# ============================================================
add_text(slide, 24.2, 13.3, 10.5, 9.2, [
    ("\u2463 Conclusions & Impact", 22, B, TITLE_COLOR),
    ("", 6, False, BODY_COLOR),
    ("Key Insight", 16, B, GREEN),
    ("Every method discovers coordination.", 14, False, BODY_COLOR),
    ("The challenge is keeping it.", 14, B, BODY_COLOR),
    ("VABL\u2019s prediction loss prevents belief drift,", 14, False, BODY_COLOR),
    ("sustaining coordination without communication.", 14, False, BODY_COLOR),
    ("", 8, False, BODY_COLOR),
    ("Data Science Contributions", 16, B, TITLE_COLOR),
    ("\u2022 Prediction from behavioral time-series: each agent", 12, False, SUB_COLOR),
    ("  processes sequential teammate action data to build", 12, False, SUB_COLOR),
    ("  a real-time predictive model (86% accuracy)", 12, False, SUB_COLOR),
    ("\u2022 Self-supervised regularization: auxiliary prediction", 12, False, SUB_COLOR),
    ("  task calibrates learned representations, preventing", 12, False, SUB_COLOR),
    ("  catastrophic forgetting (2.3\u00d7 variance reduction)", 12, False, SUB_COLOR),
    ("\u2022 Rigorous evaluation: 5 seeds, ablation studies,", 12, False, SUB_COLOR),
    ("  4 baselines, 3 environments, 10M-step extended runs", 12, False, SUB_COLOR),
    ("", 8, False, BODY_COLOR),
    ("Real-World Applications", 16, B, TITLE_COLOR),
    ("\u2022 Warehouse robot fleets (comm. blackouts)", 13, False, SUB_COLOR),
    ("\u2022 Autonomous vehicle coordination", 13, False, SUB_COLOR),
    ("\u2022 Disaster response drone swarms", 13, False, SUB_COLOR),
    ("\u2022 Collaborative surgical robotics", 13, False, SUB_COLOR),
    ("", 8, False, BODY_COLOR),
    ("Future Work", 14, B, TITLE_COLOR),
    ("\u2022 Hanabi evaluation (card-game partial observability)", 12, False, GRAY),
    ("\u2022 Scaling to N>10 agents with sparse attention", 12, False, GRAY),
    ("\u2022 Hybrid implicit + explicit communication", 12, False, GRAY),
])


# ============================================================
# Remove other template slides (keep only slide 4)
# ============================================================
for idx in sorted([0, 1, 2, 4], reverse=True):
    rId = prs.slides._sldIdLst[idx].rId
    prs.part.drop_rel(rId)
    del prs.slides._sldIdLst[idx]

prs.save('poster/VABL_Poster_v3.pptx')
print("Poster saved: poster/VABL_Poster_v2.pptx")

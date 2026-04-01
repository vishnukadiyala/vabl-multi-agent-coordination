"""VABL DISC Poster — Final build with exact template alignment.

Template Slide 4 geometry (36x24 inches):
- Title bar: full width, 0-2.9"
- 5 panels, each with crimson header (0.86" tall, #870002) + content area
  Panel A (top-left):     (0.71, 3.27)  11.39 x 9.55  — content starts at y=4.13
  Panel B (top-right):    (12.32, 3.27) 22.93 x 9.55  — content starts at y=4.13
  Panel C (bot-left):     (0.73, 13.11) 11.39 x 9.55  — content starts at y=13.97
  Panel D (bot-center):   (12.34, 13.11) 11.39 x 9.55 — content starts at y=13.97
  Panel E (bot-right):    (23.94, 13.11) 11.39 x 9.55 — content starts at y=13.97
- Footer line at y=23.1, footer text at y=23.2
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from PIL import Image

# Constants
INSET = 0.3  # padding inside panels
CRIMSON = RGBColor(0x87, 0x00, 0x02)
DARK = RGBColor(0x1A, 0x1A, 0x1A)
BODY = RGBColor(0x33, 0x33, 0x33)
SUB = RGBColor(0x66, 0x66, 0x66)
GREEN = RGBColor(0x2E, 0x7D, 0x32)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GREEN_BG = RGBColor(0xF0, 0xF7, 0xF0)

# Panel content areas (after header bar)
PANELS = {
    'A': {'x': 0.71, 'y': 4.13, 'w': 11.39, 'h': 8.69},   # top-left
    'B': {'x': 12.32, 'y': 4.13, 'w': 22.93, 'h': 8.69},   # top-right (wide)
    'C': {'x': 0.73, 'y': 13.97, 'w': 11.39, 'h': 8.69},   # bot-left
    'D': {'x': 12.34, 'y': 13.97, 'w': 11.39, 'h': 8.69},  # bot-center
    'E': {'x': 23.94, 'y': 13.97, 'w': 11.39, 'h': 8.69},  # bot-right
}

prs = Presentation('poster/Poster Templates.pptx')
slide = prs.slides[3]


def cx(panel, offset=0):
    """Content x position with inset."""
    return PANELS[panel]['x'] + INSET + offset

def cy(panel, offset=0):
    """Content y position with inset."""
    return PANELS[panel]['y'] + INSET + offset

def cw(panel, margin=0):
    """Content width with insets."""
    return PANELS[panel]['w'] - 2 * INSET - margin

def ch(panel, margin=0):
    """Content height with insets."""
    return PANELS[panel]['h'] - 2 * INSET - margin


def add_text_block(slide, x, y, w, h, lines):
    """Add a text block. Lines are (text, size_pt, bold, color) or strings."""
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if isinstance(line, str):
            line = (line, 15, False, BODY)
        text, size, bold, color = line
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if text == '':
            p.space_after = Pt(6)
            continue
        run = p.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
    return txBox


def add_figure(slide, path, x, y, max_w, max_h):
    """Add figure maintaining aspect ratio within max bounds."""
    img = Image.open(path)
    ratio = img.width / img.height
    # Fit within bounds
    w = max_w
    h = w / ratio
    if h > max_h:
        h = max_h
        w = h * ratio
    return slide.shapes.add_picture(path, Inches(x), Inches(y), Inches(w), Inches(h))


# ================================================================
# TITLE BAR — update existing shapes
# ================================================================
for shape in slide.shapes:
    if not shape.has_text_frame:
        continue
    t = shape.text_frame.text
    if 'Project Title' in t:
        shape.text_frame.clear()
        p = shape.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = 'Learning to Coordinate Without Talking'
        run.font.size = Pt(32)
        run.font.bold = True
        run.font.color.rgb = WHITE

        p2 = shape.text_frame.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        run2 = p2.add_run()
        run2.text = 'Data-Driven Belief Models for Robust Multi-Agent Coordination'
        run2.font.size = Pt(20)
        run2.font.color.rgb = RGBColor(0xCC, 0xEE, 0xCC)

        p3 = shape.text_frame.add_paragraph()
        p3.alignment = PP_ALIGN.CENTER
        run3 = p3.add_run()
        run3.text = 'Vishnu Kadiyala, Mohammed Atiquzzaman \u2014 School of Computer Science, University of Oklahoma'
        run3.font.size = Pt(14)
        run3.font.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)

    elif 'Logo' in t:
        shape.text_frame.clear()
        p = shape.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = 'OU'
        run.font.size = Pt(28)
        run.font.bold = True
        run.font.color.rgb = WHITE

    elif 'Acknowledgement' in t:
        shape.text_frame.clear()
        p = shape.text_frame.paragraphs[0]
        run = p.add_run()
        run.text = 'vishnupk@ou.edu  \u2502  github.com/vishnukadiyala/vabl-multi-agent-coordination'
        run.font.size = Pt(11)
        run.font.color.rgb = SUB


# ================================================================
# UPDATE PANEL HEADERS — modify the group header text
# ================================================================
header_names = {
    'Group 22': 'THE PROBLEM',
    'Group 24': 'VABL: HOW IT WORKS',
    'Group 31': 'EVIDENCE',
    'Group 36': 'RESULTS',
    'Group 39': 'CONCLUSIONS & IMPACT',
}
for shape in slide.shapes:
    if shape.shape_type == 6 and shape.name in header_names:
        for child in shape.shapes:
            if child.has_text_frame and 'SECTION NAME' in child.text_frame.text:
                child.text_frame.clear()
                p = child.text_frame.paragraphs[0]
                p.alignment = PP_ALIGN.CENTER
                run = p.add_run()
                run.text = header_names[shape.name]
                run.font.size = Pt(20)
                run.font.bold = True
                run.font.color.rgb = WHITE


# ================================================================
# PANEL A — THE PROBLEM (top-left)
# ================================================================
add_text_block(slide, cx('A'), cy('A'), cw('A'), 4.5, [
    ('Can AI agents coordinate', 20, True, DARK),
    ('without communication?', 20, True, DARK),
    ('', 6, False, BODY),
    ('A QB reads his receiver and throws \u2014 no', 14, False, BODY),
    ('audible needed. AI agents should coordinate', 14, False, BODY),
    ('by predicting teammates, not communicating.', 14, False, BODY),
    ('', 4, False, BODY),
    ('Setting: Dec-POMDPs \u2014 partial observations,', 13, False, SUB),
    ('no shared communication channel.', 13, False, SUB),
    ('', 6, False, BODY),
    ('Coordination Collapse', 16, True, CRIMSON),
    ('Agents learn teamwork then forget it \u2014 like', 13, False, BODY),
    ('a team falling apart in the 4th quarter.', 13, False, BODY),
])

# Collapse figure — larger, fills bottom of Panel A
add_figure(slide, 'poster/fig_collapse_all_methods.png',
    cx('A'), cy('A', 4.5), cw('A'), ch('A') - 4.5)


# ================================================================
# PANEL B — HOW VABL WORKS (top-right, wide)
# ================================================================
# Split wide panel: left half = architecture + method, right half = hero figure

half_w = cw('B') / 2 - 0.2  # small gap between halves

# Architecture diagram (left half, top)
add_figure(slide, 'poster/fig_architecture_v2.png',
    cx('B'), cy('B'), half_w, 3.5)

# Method text (left half, below architecture)
add_text_block(slide, cx('B'), cy('B', 3.7), half_w, 4.0, [
    ('Three steps, no communication:', 15, True, DARK),
    ('\u2776 Observe teammate actions', 14, False, BODY),
    ('\u2777 Attend \u2014 MHA (h=4) weighs teammates', 14, False, BODY),
    ('\u2778 Update belief via GRU', 14, False, BODY),
    ('\u2779 Predict next actions (Theory of Mind)', 14, False, BODY),
    ('', 4, False, BODY),
    ('L = L_PPO + 0.05 \u00d7 L_prediction', 14, True, DARK),
    ('Prediction accuracy: 17% \u2192 86%', 14, True, GREEN),
    ('CTDE: centralized critic, decentralized actors', 12, False, SUB),
])

# Hero figure — collapse chart (right half)
add_figure(slide, 'poster/fig_hero_collapse.png',
    cx('B', half_w + 0.4), cy('B'), half_w, 5.5)

# TL;DR below hero
tldr = add_text_block(slide,
    cx('B', half_w + 0.4), cy('B', 5.7), half_w, 2.0, [
    ('VABL: 38% collapse, no communication.', 16, True, GREEN),
    ('MAPPO & AERIAL: 100% collapse.', 16, True, CRIMSON),
    ('TarMAC: 88% collapse, needs comm. channel.', 14, False, SUB),
])


# ================================================================
# PANEL C — EVIDENCE (bottom-left)
# ================================================================
# Variance figure fills most of this panel
add_figure(slide, 'poster/fig_variance_hero.png',
    cx('C'), cy('C'), cw('C'), 5.5)

add_text_block(slide, cx('C'), cy('C', 5.7), cw('C'), 2.5, [
    ('VABL std dev 1.7 vs MAPPO 20.3', 15, True, GREEN),
    ('across 5 seeds (Simple Coordination).', 14, False, BODY),
])


# ================================================================
# PANEL D — RESULTS (bottom-center)
# ================================================================
add_figure(slide, 'poster/fig_ablation_cramped.png',
    cx('D'), cy('D'), cw('D'), 5.0)

add_text_block(slide, cx('D'), cy('D', 5.2), cw('D'), 3.2, [
    ('Key Results (5 seeds, 95% CIs)', 16, True, DARK),
    ('', 4, False, BODY),
    ('\u2022 Overcooked: 3.3\u00d7 higher reward than MAPPO', 13, False, BODY),
    ('\u2022 38% collapse vs 100% (MAPPO, AERIAL)', 13, False, BODY),
    ('\u2022 TarMAC (comm.): Final 56.6 but 88% collapse', 13, False, BODY),
    ('\u2022 5-agent: VABL 95.7\u00b13.3 vs MAPPO 84.0\u00b110.4', 13, False, BODY),
    ('\u2022 10M steps: MAPPO collapses to zero reward', 13, False, BODY),
    ('', 4, False, BODY),
    ('Cramped Room ablation confirms: Full VABL =', 13, True, GREEN),
    ('highest peak (1030\u00b170), lowest variance.', 13, True, GREEN),
])


# ================================================================
# PANEL E — CONCLUSIONS & IMPACT (bottom-right)
# ================================================================
add_text_block(slide, cx('E'), cy('E'), cw('E'), ch('E'), [
    ('Key Insight', 18, True, GREEN),
    ('', 4, False, BODY),
    ('Every method discovers coordination.', 15, False, DARK),
    ('The challenge is keeping it.', 15, True, DARK),
    ('', 4, False, BODY),
    ('VABL\u2019s auxiliary prediction loss prevents', 14, False, BODY),
    ('belief drift, sustaining coordination', 14, False, BODY),
    ('without any communication channel.', 14, False, BODY),
    ('', 10, False, BODY),
    ('Data Science Contributions', 16, True, CRIMSON),
    ('', 4, False, BODY),
    ('\u2022 Behavioral prediction from sequential', 13, False, BODY),
    ('  observational data (86% accuracy)', 13, False, SUB),
    ('\u2022 Self-supervised regularization prevents', 13, False, BODY),
    ('  catastrophic forgetting (2.3\u00d7 variance \u2193)', 13, False, SUB),
    ('\u2022 Rigorous evaluation: 5 seeds, 4 baselines,', 13, False, BODY),
    ('  3 environments, 10M-step extended runs', 13, False, SUB),
    ('', 10, False, BODY),
    ('Applications', 16, True, CRIMSON),
    ('', 4, False, BODY),
    ('\u2022 Warehouse robot fleets', 14, False, BODY),
    ('\u2022 Autonomous vehicle coordination', 14, False, BODY),
    ('\u2022 Disaster response drone swarms', 14, False, BODY),
    ('', 10, False, BODY),
    ('Future Work', 14, True, CRIMSON),
    ('\u2022 Hanabi, scaling to N>10, hybrid communication', 12, False, SUB),
])


# ================================================================
# Remove other template slides
# ================================================================
for idx in sorted([0, 1, 2, 4], reverse=True):
    rId = prs.slides._sldIdLst[idx].rId
    prs.part.drop_rel(rId)
    del prs.slides._sldIdLst[idx]

prs.save('poster/VABL_Poster_FINAL_v3.pptx')
print('Poster saved: poster/VABL_Poster_FINAL.pptx')

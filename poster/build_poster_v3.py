"""Build VABL DISC poster v3 — professionally designed layout.

Uses exact template panel coordinates from Slide 4 of Poster Templates.pptx.
All positions respect panel borders with consistent 0.35" inner padding.

Template panel geometry (measured from .pptx):
  Title bar:    (0.000, 0.000)  36.000 x 2.852
  P1 top-left:  (0.712, 3.268)  11.388 x 9.552  -> to (12.100, 12.820)
  P2 top-right: (12.317, 3.268) 22.926 x 9.552  -> to (35.243, 12.820)
  P3 bot-left:  (0.733, 13.109) 11.388 x 9.552  -> to (12.122, 22.661)
  P4 bot-center:(12.338, 13.109)11.388 x 9.552  -> to (23.727, 22.661)
  P5 bot-right: (23.943, 13.109)11.388 x 9.552  -> to (35.331, 22.661)

Each panel has a 0.86" crimson header bar at top.
Content starts at panel_top + 0.86 + 0.30 = panel_top + 1.16
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# =================================================================
# Color palette — harmonizes with template crimson #870002
# =================================================================
CRIMSON = RGBColor(0x87, 0x00, 0x02)       # Template's own crimson
DARK_TEXT = RGBColor(0x2D, 0x2D, 0x2D)      # Primary body text
MED_TEXT = RGBColor(0x40, 0x40, 0x40)        # Secondary body text
LIGHT_TEXT = RGBColor(0x80, 0x80, 0x80)      # Tertiary / footnotes
VABL_GREEN = RGBColor(0x2E, 0x7D, 0x32)     # VABL emphasis
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
SUBTITLE_PINK = RGBColor(0xFF, 0xD5, 0xD5)  # Harmonizes with crimson
AUTHOR_GRAY = RGBColor(0xE0, 0xE0, 0xE0)

B = True  # shorthand for bold


def add_text(slide, left, top, width, height, lines, default_size=15):
    """Add a text box with formatted lines.

    Each line is either a plain string or a tuple:
        (text, font_size_pt, bold, color)
    An empty string "" inserts a spacer paragraph.
    """
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top),
                                     Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, line in enumerate(lines):
        if isinstance(line, str):
            line = (line, default_size, False, MED_TEXT)
        text, size, bold, color = line

        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()

        if text == "":
            p.space_after = Pt(6)
            continue

        run = p.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color


# =================================================================
# Load template
# =================================================================
prs = Presentation('poster/Poster Templates.pptx')
slide = prs.slides[3]  # Slide 4 template


# =================================================================
# TITLE BAR — use existing template shapes
# =================================================================
for shape in slide.shapes:
    if not shape.has_text_frame:
        continue
    text = shape.text_frame.text

    if 'Project Title' in text:
        shape.text_frame.clear()

        # Title line
        p = shape.text_frame.paragraphs[0]
        run = p.add_run()
        run.text = "Learning to Coordinate Without Talking"
        run.font.size = Pt(36)
        run.font.bold = True
        run.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

        # Subtitle
        p2 = shape.text_frame.add_paragraph()
        run2 = p2.add_run()
        run2.text = "Data-Driven Belief Models for Robust Multi-Agent Coordination"
        run2.font.size = Pt(24)
        run2.font.color.rgb = SUBTITLE_PINK
        p2.alignment = PP_ALIGN.CENTER

        # Authors
        p3 = shape.text_frame.add_paragraph()
        run3 = p3.add_run()
        run3.text = ("Vishnu Kadiyala, Mohammed Atiquzzaman  \u2014  "
                     "School of Computer Science, University of Oklahoma")
        run3.font.size = Pt(16)
        run3.font.color.rgb = AUTHOR_GRAY
        p3.alignment = PP_ALIGN.CENTER

    elif 'Logo' in text:
        shape.text_frame.clear()
        p = shape.text_frame.paragraphs[0]
        run = p.add_run()
        run.text = "OU"
        run.font.size = Pt(28)
        run.font.bold = True
        run.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

    elif 'Acknowledgement' in text:
        shape.text_frame.clear()
        p = shape.text_frame.paragraphs[0]
        run = p.add_run()
        run.text = ("Contact: vishnupk@ou.edu   \u2502   "
                    "github.com/vishnukadiyala/vabl-multi-agent-coordination   "
                    "\u2502   DISC 2026")
        run.font.size = Pt(13)
        run.font.color.rgb = LIGHT_TEXT


# =================================================================
# PANEL 1: THE PROBLEM (top-left)
# Panel: (0.712, 3.268) 11.388 x 9.552
# Content area: (1.06, 4.43) to (11.75, 12.47)
# =================================================================

# Hook question
add_text(slide, 1.06, 4.43, 10.39, 1.70, [
    ("Can AI agents learn to coordinate", 20, B, DARK_TEXT),
    ("without communicating?", 20, B, DARK_TEXT),
])

# Football analogy
add_text(slide, 1.06, 6.20, 10.39, 1.60, [
    ("A quarterback reads the safety, sees his receiver", 15, False, MED_TEXT),
    ("break right, and throws \u2014 no audible needed.", 15, False, MED_TEXT),
    ("He coordinates by predicting, not communicating.", 15, B, MED_TEXT),
    ("VABL teaches AI agents the same skill.", 15, B, VABL_GREEN),
])

# Collapse subheading
add_text(slide, 1.06, 7.90, 10.39, 0.45, [
    ("The Hidden Problem: Coordination Collapse", 18, B, CRIMSON),
])

# Collapse figure — respects panel width
# fig_collapse_all_methods.png: aspect 2.60, so 10.19" wide -> 3.92" tall
slide.shapes.add_picture('poster/fig_collapse_all_methods.png',
    Inches(1.16), Inches(8.45), Inches(10.19), Inches(3.92))


# =================================================================
# PANEL 2: HOW VABL WORKS (top-right, wide)
# Panel: (12.317, 3.268) 22.926 x 9.552
# Content area: (12.67, 4.43) to (34.89, 12.47)
# Split: LEFT half 12.67-23.17 (10.5"), RIGHT half 23.67-34.89 (11.22")
# =================================================================

# --- LEFT HALF: Architecture + Method Steps ---

# Architecture diagram
# fig_architecture_v2.png: aspect 2.35, so 10.50" wide -> 4.47" tall
slide.shapes.add_picture('poster/fig_architecture_v2.png',
    Inches(12.67), Inches(4.43), Inches(10.50), Inches(4.47))

# Method steps (below architecture)
add_text(slide, 12.67, 9.10, 10.50, 3.20, [
    ("\u2776 Observe \u2014 Encode teammate actions via MLP", 15, B, CRIMSON),
    ("    One-hot action \u2192 64-dim embedding per teammate", 13, False, MED_TEXT),
    ("\u2777 Attend \u2014 Multi-head attention (h=4) over actions", 15, B, CRIMSON),
    ("    Query = belief state, Keys/Values = action embeddings", 13, False, MED_TEXT),
    ("\u2778 Update \u2014 GRU: b_t = GRU([obs || context], b_{t\u22121})", 15, B, CRIMSON),
    ("\u2779 Predict \u2014 Auxiliary MLP predicts next actions", 15, B, CRIMSON),
    ("    Accuracy: 17% (chance) \u2192 86%. Keeps beliefs calibrated.", 13, False, MED_TEXT),
    ("", 6, False, MED_TEXT),
    ("L_total = L_PPO + 0.05 \u00d7 L_aux     |     <10% overhead vs MAPPO", 14, B, DARK_TEXT),
])

# --- RIGHT HALF: Hero Figure + TL;DR ---

# Hero figure — THE showstopper, largest figure on poster
# fig_hero_retained.png: aspect 1.70, so 10.72" wide -> 6.31" tall
slide.shapes.add_picture('poster/fig_hero_retained.png',
    Inches(23.67), Inches(4.43), Inches(10.72), Inches(6.31))

# TL;DR callout box (below hero figure)
# Positioned to leave breathing room below the hero
tldr = slide.shapes.add_textbox(
    Inches(23.67), Inches(11.00), Inches(10.72), Inches(1.30))
tf = tldr.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
run = p.add_run()
run.text = ("VABL retains 62% of coordination vs. 0% for MAPPO "
            "and AERIAL \u2014 without any communication channel.")
run.font.size = Pt(16)
run.font.bold = True
run.font.color.rgb = VABL_GREEN
p.alignment = PP_ALIGN.CENTER
# Note: manually add #F0F7F0 fill and #2E7D32 border in PowerPoint


# =================================================================
# PANEL 3: EVIDENCE — RELIABILITY (bottom-left)
# Panel: (0.733, 13.109) 11.388 x 9.552
# Content area: (1.08, 14.27) to (11.77, 22.31)
# =================================================================

# Variance figure — one big figure per panel
# fig_variance_v2.png: aspect 1.64, so 10.29" wide -> 6.27" tall
slide.shapes.add_picture('poster/fig_variance_v2.png',
    Inches(1.18), Inches(14.27), Inches(10.29), Inches(6.27))

# Key stats below figure
add_text(slide, 1.18, 20.70, 10.29, 1.50, [
    ("\u2022 VABL std: 1.7 vs MAPPO: 20.3 (12\u00d7 more reliable)", 14, B, VABL_GREEN),
    ("\u2022 Overcooked AA: VABL 3.3\u00d7 MAPPO final reward", 14, False, MED_TEXT),
    ("\u2022 5 seeds, 95% CIs, extended 10M-step runs", 14, False, MED_TEXT),
])


# =================================================================
# PANEL 4: RESULTS SUMMARY (bottom-center)
# Panel: (12.338, 13.109) 11.388 x 9.552
# Content area: (12.69, 14.27) to (23.38, 22.31)
# =================================================================

# Results comparison (text-based table — replace manually with real table)
add_text(slide, 12.89, 14.37, 10.29, 4.80, [
    ("Head-to-Head Comparison", 18, B, CRIMSON),
    ("", 6, False, MED_TEXT),
    ("                     VABL    MAPPO   TarMAC  AERIAL", 13, B, DARK_TEXT),
    ("\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500", 10, False, LIGHT_TEXT),
    ("Reward Retained    62%      0%      14%      0%", 14, False, MED_TEXT),
    ("Collapse %          38%    100%      85%    100%", 14, False, MED_TEXT),
    ("Reward Std Dev      1.7     20.3      --       --", 14, False, MED_TEXT),
    ("Comm. Required?     No       No      Yes      No", 14, False, MED_TEXT),
    ("\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500", 10, False, LIGHT_TEXT),
    ("", 4, False, MED_TEXT),
    ("\u2192 Replace this with a proper PowerPoint table", 11, False, LIGHT_TEXT),
])

# Ablation summary
add_text(slide, 12.89, 19.40, 10.29, 2.80, [
    ("Component Ablation (Cramped Room, 5 seeds)", 16, B, CRIMSON),
    ("", 6, False, MED_TEXT),
    ("\u2022 Full VABL:     1030 \u00b1 70   (best, lowest variance)", 14, B, VABL_GREEN),
    ("\u2022 No Attention:  990 \u00b1 65", 14, False, MED_TEXT),
    ("\u2022 No Aux Loss:   951 \u00b1 162  (2.3\u00d7 higher variance)", 14, False, MED_TEXT),
    ("\u2022 Neither:       906 \u00b1 244  (3.5\u00d7 higher variance)", 14, False, MED_TEXT),
    ("", 6, False, MED_TEXT),
    ("Auxiliary prediction is the critical stabilizer.", 15, B, VABL_GREEN),
])


# =================================================================
# PANEL 5: CONCLUSIONS & IMPACT (bottom-right)
# Panel: (23.943, 13.109) 11.388 x 9.552
# Content area: (24.29, 14.27) to (34.98, 22.31)
# =================================================================

# Key Insight (accent box — add fill manually)
add_text(slide, 24.39, 14.37, 10.09, 2.10, [
    ("Key Insight", 18, B, CRIMSON),
    ("", 4, False, MED_TEXT),
    ("Every method discovers coordination.", 16, False, DARK_TEXT),
    ("The challenge is keeping it.", 18, B, CRIMSON),
    ("VABL's prediction loss prevents belief drift,", 15, False, MED_TEXT),
    ("sustaining coordination without communication.", 15, False, MED_TEXT),
])

# Data Science Contributions
add_text(slide, 24.39, 16.70, 10.09, 2.70, [
    ("Data Science Contributions", 16, B, CRIMSON),
    ("", 4, False, MED_TEXT),
    ("\u2022 Prediction from behavioral time-series:", 14, B, DARK_TEXT),
    ("  sequential teammate action data \u2192 real-time", 14, False, MED_TEXT),
    ("  predictive model (86% accuracy)", 14, False, MED_TEXT),
    ("\u2022 Self-supervised regularization prevents", 14, B, DARK_TEXT),
    ("  catastrophic forgetting (2.3\u00d7 variance reduction)", 14, False, MED_TEXT),
    ("\u2022 Rigorous evaluation: 5 seeds, 4 baselines,", 14, False, MED_TEXT),
    ("  3 environments, 10M-step extended runs", 14, False, MED_TEXT),
])

# Applications
add_text(slide, 24.39, 19.60, 10.09, 1.50, [
    ("Real-World Applications", 16, B, CRIMSON),
    ("\u2022 Warehouse robot fleets (comm. blackouts)", 14, False, MED_TEXT),
    ("\u2022 Autonomous vehicle coordination", 14, False, MED_TEXT),
    ("\u2022 Disaster response drone swarms", 14, False, MED_TEXT),
])

# Future Work
add_text(slide, 24.39, 21.20, 10.09, 1.00, [
    ("Future Work", 14, B, CRIMSON),
    ("\u2022 Hanabi evaluation  \u2022  Scale to N>10 agents  \u2022  Hybrid communication", 13, False, LIGHT_TEXT),
])


# =================================================================
# Remove other template slides (keep only slide 4)
# =================================================================
for idx in sorted([0, 1, 2, 4], reverse=True):
    rId = prs.slides._sldIdLst[idx].rId
    prs.part.drop_rel(rId)
    del prs.slides._sldIdLst[idx]

prs.save('poster/VABL_Poster_v3_pro.pptx')
print("Poster saved: poster/VABL_Poster_v3_pro.pptx")
print()
print("NEXT STEPS — manual polish in PowerPoint:")
print("  1. Enter each panel group -> replace 'SECTION NAME' with proper titles")
print("  2. Replace 'OU' text with actual OU logo image")
print("  3. Add #F0F7F0 fill + #2E7D32 border to TL;DR box (Panel 2)")
print("  4. Add #FFF8E1 fill + crimson left accent bar to Key Insight (Panel 5)")
print("  5. Replace text table in Panel 4 with a real PowerPoint table")
print("  6. Add 0.75pt #D0D0D0 borders to all figures")
print("  7. Generate QR code and place at (0.90, 23.20) in footer")
